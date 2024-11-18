import json, os, re, logging, copy, traceback
import httpx, asyncio, tiktoken, fitz
import pandas as pd
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation

from jupyter_server.base.handlers import APIHandler
from jupyter_server.utils import url_path_join
from concurrent_log_handler import ConcurrentRotatingFileHandler
import tornado

# 创建一个名为 'notechat_logger' 的专用日志记录器
notechat_logger = logging.getLogger("notechat_logger")
notechat_logger.setLevel(logging.INFO)

# 检查@log文件夹是否存在，如果不存在则创建
log_directory = "@log"
if not os.path.exists(log_directory):
    os.makedirs(log_directory)

# 加载logging配置
notechat_fh = ConcurrentRotatingFileHandler(
    os.path.join(log_directory, "notechat.log"),
    maxBytes=10*1024*1024,
    backupCount=100,
    encoding='utf-8'
)
notechat_fh.setFormatter(logging.Formatter(fmt="%(asctime)s - %(levelname)s - %(message)s", datefmt="%Y-%m-%d %H:%M:%S"))
notechat_logger.addHandler(notechat_fh)
notechat_logger.info(f"###### IT IS A GOOD DAY TODAY, LET'S FIND 42 !! ######")


class RouteHandler(APIHandler):
    # The following decorator should be present on all verb methods (head, get, post,
    # patch, put, delete, options) to ensure only authorized user can request the
    # Jupyter server
    @tornado.web.authenticated
    def get(self):
        notechat_logger.info(f"###### TEST ENDPOINT: /jupyterlab-notechat/get-example ######")
        self.finish(json.dumps({
            "data": "This is /jupyterlab-notechat/get-example endpoint! You are authenticated!"
        }))


class ChatHandler(APIHandler):
    @tornado.web.authenticated
    async def post(self):       
        try:
            # 从请求体中获取messages
            data = json.loads(self.request.body)
            cell_json_arr = data.get("cell_json_arr", [])
            active_cell_index = data.get("active_cell_index", 0)
            ai_name = data.get("ai_name", "**assistant**")
            user_name = data.get("user_name", "**user**")
            ref_name = data.get("ref_name", "_ref")
            prompt = (str(data.get("prompt", "You are a helpful and warm-hearted assistant:)")) + " " + data.get("add_prompt", "")).strip()
            model = data.get("model", "gpt-4o")
            vision_model = data.get("vision_model", "gpt-4o")
            # 目前只支持openai的vision模型，如果使用moonshot则不使用vision模型
            use_vision = False if model.startswith("moonshot") else parse_param(data, "use_vision", bool, True)
            max_input = parse_param(data, "max_input", int, 80000) # data.get("max_input", 80000)
            # 如果max_output为0，则设置为None，表示不限制
            max_output = parse_param(data, "max_output", int, 0) # data.get("max_output", 0)
            max_output = None if max_output <= 0 else max_output
            temperature = parse_param(data, "temperature", float, 0.6) # data.get("temperature", 0.6)
            response_format = data.get("response_format", "text")
            timeout = parse_param(data, "timeout", int, 600) # timeout = data.get("timeout", 600)
            retries = parse_param(data, "retries", int, 3) # retries = data.get("retries", 3)
            delay = parse_param(data, "delay", int, 1) # delay = data.get("delay", 1)
            api_key = data.get("openai_api_key", "None")
            moonshot_api_key = data.get("moonshot_api_key", "None")
            files = data.get("files", "").split(" ")
            
            # 从cell_json_arr，system prompt和files生成messages
            messages, has_image, total_input, input_tokens = await self.get_all_messages(cell_json_arr, active_cell_index, ai_name, user_name, ref_name, model, use_vision, max_input, prompt, files)

            # 如果model是moonshot，调用moonshot_chat函数
            if model.startswith("moonshot"):
                notechat_logger.info(f"### PARAMS ### model: {model} || use_vision: {use_vision} || has_image: {has_image} || max_input: {max_input} || total_input: {total_input} || input_tokens: {input_tokens} || max_output: {max_output}  || temperature: {temperature} || files: {files} ||")
                notechat_logger.info(f"### INPUT MESSAGES ### {messages}")
                response = await self.moonshot_chat(messages, model, max_output, temperature, timeout, retries, delay, moonshot_api_key)
            # 其他情况，调用openai_chat函数
            else:
                if has_image and use_vision:
                    notechat_logger.info(f"### PARAMS ### model: {vision_model} || use_vision: {use_vision} || has_image: {has_image} || max_input: {max_input} || total_input: {total_input} || input_tokens: {input_tokens} || max_output: {max_output}  || temperature: {temperature} || files: {files} ||")
                    logging_messages = copy.deepcopy(messages)
                    for message in logging_messages:
                        if isinstance(message["content"], list):
                            for content in message["content"]:
                                if content["type"] == "image_url":
                                    content["image_url"]["url"] = content["image_url"]["url"][0:30] + "..." + content["image_url"]["url"][-30:]
                    notechat_logger.info(f"### INPUT MESSAGES ### {logging_messages}")
                    response = await self.openai_chat(messages, vision_model, max_output, None, temperature, timeout, retries, delay, api_key)
                else:
                    notechat_logger.info(f"### PARAMS ### model: {model} || use_vision: {use_vision} || has_image: {has_image} || max_input: {max_input} || total_input: {total_input} || input_tokens: {input_tokens} || max_output: {max_output}  || temperature: {temperature} || files: {files} ||")
                    notechat_logger.info(f"### INPUT MESSAGES ### {messages}")
                    response = await self.openai_chat(messages, model, max_output, response_format, temperature, timeout, retries, delay, api_key)

            notechat_logger.info(f"### OUTPUT RESPONSE ### {response}")

            self.finish(json.dumps(response))

        except Exception as e:
            # 设置HTTP状态码为500（内部服务器错误）
            self.set_status(500)
            self.finish(json.dumps({"error": "API请求处理出错: " + str(e)}))

    async def openai_chat(self, messages, model="gpt-4o", max_tokens=None, response_format="text", temperature=0.6, timeout=300, retries=3, delay=1, api_key=None):
        """
        使用OpenAI API进行对话生成

        Args:
            messages: 对话消息列表
            model: 模型名称，gpt-4o，o1-preview
            max_tokens: 最大生成长度
            response_format: 响应格式，值为`text`、`json_object`、None
            temperature: 温度参数
            timeout: 超时秒数
            retries: 重试次数
            delay: 重试延迟秒数
        """
        # 首先检查环境变量中的 OPENAI_API_KEY
        env_api_key = os.environ.get("OPENAI_API_KEY")
        if env_api_key:
            # 如果环境变量中的 OPENAI_API_KEY 存在且非空，优先使用环境变量中的值
            api_key = env_api_key
        elif api_key is None or api_key.lower() == "none":
            # 如果传入的 api_key 不存在或其值为 "none"（不区分大小写）
            return {"error": 'OpenAI API Key Missing ... 2 ways to setup api keys: 1. Top Menu Bar -> Settings -> Settings Editor -> NoteChat -> Param `@openai_api_key` ; 2. set key to server environment variable `OPENAI_API_KEY`, linux `export OPENAI_API_KEY=your_key`, windows `$env:OPENAI_API_KEY = "your_key"`'}
               
        url = "https://api.openai.com/v1/chat/completions"
        headers = {
            "Content-Type": "application/json",
            "Authorization": "Bearer " + api_key,
        }
        data = {
            "model": model,
            "messages": messages,
            "temperature": temperature,
        }
        if max_tokens is not None:
            data["max_tokens"] = max_tokens
        if response_format is not None:
            data["response_format"] = {"type": response_format}
        # 如果是o1模型，则有一些参数要特殊处理
        if model.startswith("o1"):
            # 检查message的第一个message的role是否是system，如果是则替换成user
            if len(messages)>0 and messages[0]["role"] == "system":
                messages[0]["role"] = "user"
            # 检查是否有max_tokens参数，如果有则替换成max_completion_tokens
            if "max_tokens" in data:
                data["max_completion_tokens"] = data["max_tokens"]
                del data["max_tokens"]
            # 检查temperature如果不是1，则强制设定为1
            if temperature != 1.0:
                data["temperature"] = 1.0
        # proxy = os.environ["http_proxy"]
        # async with httpx.AsyncClient(proxies={"http://": "http://"+proxy, "https://": "https://"+proxy}) as client:
        async with httpx.AsyncClient() as client:
            attempt = 0
            while attempt < retries:
                try:
                    response = await client.post(url, headers=headers, json=data, timeout=timeout)
                    return response.json()
                except Exception as e:
                    logging.error(f"尝试 {attempt+1} / {retries}: 错误 - {str(e)}")
                    await asyncio.sleep(delay)
                    attempt += 1

        # 在达到最大重试次数后，返回错误信息，而不是抛出异常
        return {"error": f"API请求失败，已重试{retries}次，无法获取响应"}
    
    async def moonshot_chat(self, messages, model="moonshot-v1-8k", max_tokens=None, temperature=0.3, timeout=300, retries=3, delay=1, api_key=None):
        """
        使用MOONSHOT API进行对话生成

        Args:
            messages: 对话消息列表
            model: 模型名称，moonshot-v1-8k，moonshot-v1-32k，moonshot-v1-128k
            max_tokens: 最大生成长度
            temperature: 温度参数
            timeout: 超时秒数
            retries: 重试次数
            delay: 重试延迟秒数
        """
        # 首先检查环境变量中的 MOONSHOT_API_KEY
        env_api_key = os.environ.get("MOONSHOT_API_KEY")
        if env_api_key:
            # 如果环境变量中的 OPENAI_API_KEY 存在且非空，优先使用环境变量中的值
            api_key = env_api_key
        elif api_key is None or api_key.lower() == "none":
            # 如果传入的 api_key 不存在或其值为 "none"（不区分大小写）
            return {"error": 'KIMI MOONSHOT API Key Missing ... 2 ways to setup api keys: 1. Top Menu Bar -> Settings -> Settings Editor -> NoteChat -> Param `@moonshot_api_key` ; 2. set key to server environment variable `MOONSHOT_API_KEY`, linux `export MOONSHOT_API_KEY=your_key`, windows `$env:MOONSHOT_API_KEY = "your_key"`'}
        
        url = "https://api.moonshot.cn/v1/chat/completions"
        headers = {
            "Content-Type": "application/json",
            "Authorization": "Bearer " + api_key,
        }
        data = {
            "model": model,
            "messages": messages,
            "temperature": temperature,
        }
        if max_tokens is not None:
            data["max_tokens"] = max_tokens
        
        async with httpx.AsyncClient() as client:
            attempt = 0
            while attempt < retries:
                try:
                    response = await client.post(url, headers=headers, json=data, timeout=timeout)
                    return response.json()
                except Exception as e:
                    logging.error(f"尝试 {attempt+1} / {retries}: 错误 - {str(e)}")
                    await asyncio.sleep(delay)
                    attempt += 1

        # 在达到最大重试次数后，返回错误信息，而不是抛出异常
        return {"error": f"API请求失败，已重试{retries}次，无法获取响应"}
    
    # 将当前notebook指定的cell和指定的files生成消息
    async def get_all_messages(self, cell_json_arr, active_cell_index, ai_name, user_name, ref_name, model, use_vision, max_input, prompt, files):

        # 根据files生成引用文件的所有的messasges
        try:
            file_messages, file_tokens, file_has_image = await files_to_message(files, ai_name, user_name, ref_name, model, use_vision)
        except Exception as e:
            error_traceback = traceback.format_exc()  # 获取完整的堆栈跟踪
            content_text = f"########Fetching files `{files}` content error########\n\nException:\n{e}\n\nTraceBack:\n{error_traceback}"
            file_messages = [{
                "role": "user",
                "name": "context",
                "content": content_text
            }]
            file_tokens = [get_num_tokens(content_text, model)]
            file_has_image = False
            notechat_logger.error(f"### FILES PROCESSING ERROR ### Exception: {str(e)} || TraceBack: {error_traceback} ||")
        
        # 生成本notebook传回的所有messages
        try:
            notebook_messages, notebook_tokens, notebook_has_image = await cell_json_to_message(cell_json_arr, active_cell_index, ai_name, user_name, ref_name, model, use_vision)
        except Exception as e:
            error_traceback = traceback.format_exc()  # 获取完整的堆栈跟踪
            content_text = f"########Fetching current notebook content error########\n\nException:\n{e}\n\nTraceBack:\n{error_traceback}"
            notebook_messages = [{
                "role": "user",
                "name": "context",
                "content": content_text
            }]
            notebook_tokens = [get_num_tokens(content_text, model)]
            notebook_has_image = False
            notechat_logger.error(f"### NOTEBOOK PROCESSING ERROR ### Exception: {str(e)} || TraceBack: {error_traceback} ||")

        # 如果有跨文件引用，则为当前notebook的0号起始位置插入提示文件名的message用于区分文件
        # 如果没有跨文件引用，则不插入提示文件名的message，这样不用引起混淆
        if len(file_messages)>0:
            content_text = "########The following messages are generated from current working notebook########"
            notebook_messages.insert(0, {"role": "user", "name": "context", "content": content_text})
            notebook_tokens.insert(0, get_num_tokens(content_text, model))
        
        messages = file_messages + notebook_messages
        tokens = file_tokens + notebook_tokens
        has_image = file_has_image or notebook_has_image

        # 加入system message token量作为起始点，prompt肯定不为None
        prompt_token = get_num_tokens(prompt, model)
        total_input = prompt_token

        # 计算总token数量，从后向前数，如果超过则截断
        # 要注意，这里system message是tokens中的第一个id，len(messages)比len(tokens)少1
        i = len(tokens) - 1
        while i >= 0:
            total_input += tokens[i]
            if total_input > max_input:
                break
            i -= 1
        # 即便active cell是最后一个超过了字符数量，也要保证至少有一个message
        final_messages = messages[min(i + 1, len(messages)-1):]
        final_tokens = tokens[min(i + 1, len(tokens)-1):]
        # 将system message放在第一个message
        if len(prompt)>0:
            final_messages.insert(0, {"role": "system", "content": prompt})

        input_tokens = f"【{str(prompt_token)}】【{','.join(map(str, file_tokens))}】【{','.join(map(str, notebook_tokens))}】=>【{str(prompt_token)}】【{','.join(map(str, final_tokens))}】"
        return final_messages, has_image, total_input, input_tokens

# 将所有指定文件转化为消息
async def files_to_message(files, ai_name, user_name, ref_name, model, use_vision):
    messages = []
    tokens = []
    has_image = False

    for file in files:
        if file.strip() == "":
            continue
        # 读取.ipynb文件，ipynb文件特殊处理为对话型的文件，可以有多条message
        if file.endswith(".ipynb"):
            file_messages, file_tokens, file_has_image = await get_message_from_ipynb(file, ai_name, user_name, ref_name, model, use_vision)
            messages.extend(file_messages)
            tokens.extend(file_tokens)
            has_image = has_image or file_has_image
        # 其他类型的文件，都放在一个message中
        else:
            # 文本文件
            if file.endswith((".txt", ".md", ".py", ".js", ".ts", ".sh", ".bat", ".json", ".xml", ".log", ".config", ".ini", ".yaml", ".yml")):
                with open(file, "r", encoding='utf-8') as f:
                    file_text = f.read()
            # 数据文件：.csv, .xlsx, .xls
            elif file.endswith((".csv", ".xlsx", ".xls")):
                file_text = get_text_from_data(file)
            # pdf文件：.pdf
            elif file.endswith(".pdf"):
                file_text = get_text_from_pdf(file)
            # word文件：.docx
            elif file.endswith(".docx"):
                file_text = get_text_from_docx(file)
            # ppt文件：.pptx
            elif file.endswith(".pptx"):
                file_text = get_text_from_pptx(file)
            # html文件：.html, .htm
            elif file.endswith((".html", ".htm")):
                file_text = get_text_from_html(file)
            # 其他格式暂时不支持
            else:
                file_text = None

            # 解析文件名要显示出来
            if file_text is not None:
                content_text = f"########File `{file}` contains following content########\n\n{file_text}"
            # 如果文件不能解析，则在消息中提示
            else:
                content_text = f"########File format of `{file}` can not be parsed currently########"
            
            # 将文件内容放入message中
            messages.append({
                "role": "user",
                "name": "context",
                "content": content_text
            })
            tokens.append(get_num_tokens(content_text, model))

    return messages, tokens, has_image

# 将外源路径的ipynb文件转化为messages
async def get_message_from_ipynb(file_path, ai_name, user_name, ref_name, model, use_vision):

    # 打开文件
    with open(file_path, "r", encoding='utf-8') as file:
        cell_json_arr = json.load(file)["cells"]
    # 打印use vision，顺便判断是不是bool值
    messages = []
    tokens = []
    has_image = False

    # 构建用于 ai_name 的正则表达式
    ai_name_regex = re.compile(r'^{}.*\n?'.format(re.escape(ai_name)), re.IGNORECASE)

    # 构建用于 user_name 的正则表达式
    user_name_regex = re.compile(r'^{}.*\n?'.format(re.escape(user_name)), re.IGNORECASE)

    # 构建用于 div 标签的正则表达式
    ref_name_regex = re.compile(r'<div.*?>{}.*?{}.*?</div>$'.format(re.escape(ref_name), re.escape(ref_name)), re.IGNORECASE)

    for id, cell in enumerate(cell_json_arr):
        message = { 
            "role": "user",
            "name": "context"
        }
        source_text = ""
        output_text = []
        content_image = []

        cell["source"] = ''.join(cell["source"])

        # 如果source首行含有ai_name或user_name，则更换角色
        # 检查 ai_name 的匹配，如果匹配移除第一行
        if ai_name_regex.search(cell["source"]):
            cell["source"] = ai_name_regex.sub("", cell["source"])
            message["role"] = "assistant"
            message["name"] = "assistant"
        # 检查 user_name 的匹配，如果匹配移除第一行
        if user_name_regex.search(cell["source"]):
            cell["source"] = user_name_regex.sub("", cell["source"])
            message["role"] = "user"
            message["name"] = "user"

        # 如果source尾行含有ref_name，则去除尾行
        if ref_name_regex.search(cell["source"]):
            cell["source"] = ref_name_regex.sub("", cell["source"])

        # 处理source文本
        if len(cell["source"].strip())>0:
            source_text += cell["source"].strip()

        # 如果是markdown单元格，目前需要单独处理附件中的图片
        if cell["cell_type"] == "markdown":
            # 处理markdown附件                
            if "attachments" in cell:
                for _, data in cell["attachments"].items():
                    # 处理图片类型附件
                    if use_vision and "image/png" in data and len(data["image/png"])>0:
                        content_image.append({"type": "image_url", "image_url": {"url": f"data:image/png;base64," + data["image/png"] } })
                    elif use_vision and "image/jpeg" in data and len(data["image/jpeg"])>0:
                        content_image.append({"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64," + data["image/jpeg"] } })
                    elif use_vision and "image/gif" in data and len(data["image/gif"])>0:
                        content_image.append({"type": "image_url", "image_url": {"url": f"data:image/gif;base64," + data["image/gif"] } })
                    elif use_vision and "image/webp" in data and len(data["image/webp"])>0:
                        content_image.append({"type": "image_url", "image_url": {"url": f"data:image/webp;base64," + data["image/webp"] } })
                    # 目前openai vision不支持bmp格式
        
        # 如果是raw单元格，目前暂时没有特殊处理
        if cell["cell_type"] == "raw":
            pass
        
        # 如果是code单元格，目前要处理outputs中的内容
        if cell["cell_type"] == "code":
            if "outputs" in cell and len(cell["outputs"])>0:
                
                for output in cell["outputs"]:
                    # 一般是打印出来的内容
                    if output["output_type"] == "stream":
                        output["text"] = ''.join(output["text"])
                        clean_stream = remove_ansi_codes(output["text"])
                        output_text.append(clean_stream.strip())
                    # 单元格输出的错误内容
                    elif output["output_type"] == "error":
                        # 去掉traceback中的颜色类的ansi符号
                        clean_traceback = [remove_ansi_codes(text) for text in output["traceback"]]
                        clean_traceback_text = "\n".join(clean_traceback).strip()
                        output_text.append(f'''Error Name:{output["ename"]}\nError Value:{output["evalue"]}\nError Traceback:{clean_traceback_text}''')
                    elif output["output_type"] == "execute_result" or output["output_type"] == "display_data":
                        if "data" in output and len(output["data"])>0:
                            # 一般是变量输出的值
                            if "text/plain" in output["data"] and len(output["data"]["text/plain"])>=0:
                                output["data"]["text/plain"] = ''.join(output["data"]["text/plain"])
                                output_text.append(output["data"]["text/plain"].strip())
                            # 一般是plotly的微缩图片
                            if use_vision and "image/png" in output["data"] and len(output["data"]["image/png"])>0:
                                content_image.append({"type": "image_url", "image_url": {"url": f"data:image/png;base64," + output["data"]["image/png"] } })
                            elif use_vision and "image/jpeg" in output["data"] and len(output["data"]["image/jpeg"])>0:
                                content_image.append({"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64," + output["data"]["image/jpeg"] } })
                            elif use_vision and "image/gif" in output["data"] and len(output["data"]["image/gif"])>0:
                                content_image.append({"type": "image_url", "image_url": {"url": f"data:image/gif;base64," + output["data"]["image/gif"] } })
                            elif use_vision and "image/webp" in output["data"] and len(output["data"]["image/webp"])>0:
                                content_image.append({"type": "image_url", "image_url": {"url": f"data:image/webp;base64," + output["data"]["image/webp"] } })

        # 如果有图片，则标记为有图片
        if len(content_image) > 0:
            has_image = True
        
        # 准备该条信息的结构数据
        content_text = ""
        if len(source_text) > 0:
            content_text += source_text + "\n"
        if len(output_text) > 0:
            content_text += "\nexecuted outputs:\n" + "\n----------\n".join(output_text) + "\n----------"

        content_text = content_text.strip()

        if len(content_image) > 0 and len(content_text) > 0:
            message["content"] = [{"type": "text", "text": content_text}]
            message["content"].extend(content_image)
        elif len(content_image) > 0 and len(content_text) <= 0:
            message["content"] = content_image
        elif len(content_image) <= 0 and len(content_text) > 0:
            message["content"] = content_text
        else:
            continue

        messages.append(message)
        tokens.append(get_num_tokens(content_text, model))
    
    # 在0号位置插入一个message作为起始点
    content_text = f"########The following messages are generated from file `{file_path}`########"
    messages.insert(0, {"role": "user", "name": "context", "content": content_text})
    tokens.insert(0, get_num_tokens(content_text, model))

    return messages, tokens, has_image

# 将当前notebook传回的json转化为messages
async def cell_json_to_message(cell_json_arr, active_cell_index, ai_name, user_name, ref_name, model, use_vision):
    # 打印use vision，顺便判断是不是bool值
    messages = []
    tokens = []
    has_image = False
    is_active_cell_last = True
    active_cell_tokens = 0

        # 构建用于 ai_name 的正则表达式
    ai_name_regex = re.compile(r'^{}.*\n?'.format(re.escape(ai_name)), re.IGNORECASE)

    # 构建用于 user_name 的正则表达式
    user_name_regex = re.compile(r'^{}.*\n?'.format(re.escape(user_name)), re.IGNORECASE)

    # 构建用于 div 标签的正则表达式
    ref_name_regex = re.compile(r'<div.*?>{}.*?{}.*?</div>$'.format(re.escape(ref_name), re.escape(ref_name)), re.IGNORECASE)

    for id, cell in enumerate(cell_json_arr):
        message = { 
            "role": "user",
            "name": "context"
        }
        source_text = ""
        output_text = []
        content_image = []

        # 如果source首行含有ai_name或user_name，则更换角色
        # 检查 ai_name 的匹配，如果匹配移除第一行
        if ai_name_regex.search(cell["source"]):
            cell["source"] = ai_name_regex.sub("", cell["source"])
            message["role"] = "assistant"
            message["name"] = "assistant"
        # 检查 user_name 的匹配，如果匹配移除第一行
        if user_name_regex.search(cell["source"]):
            cell["source"] = user_name_regex.sub("", cell["source"])
            message["role"] = "user"
            message["name"] = "user"

        # 如果是活动单元格，强行标注为user角色
        if cell["num_id"] == active_cell_index:
            message["role"] = "user"
            message["name"] = "user"

        # 如果source尾行含有ref_name，则去除尾行
        if ref_name_regex.search(cell["source"]):
            cell["source"] = ref_name_regex.sub("", cell["source"])

        # 处理source文本
        if len(cell["source"].strip())>0:
            source_text += cell["source"].strip()

        # 如果是markdown单元格，目前需要单独处理附件中的图片
        if cell["cell_type"] == "markdown":
            # 处理markdown附件                
            if "attachments" in cell:
                for _, data in cell["attachments"].items():
                    # 处理图片类型附件
                    if use_vision and "image/png" in data and len(data["image/png"])>0:
                        content_image.append({"type": "image_url", "image_url": {"url": f"data:image/png;base64," + data["image/png"] } })
                    elif use_vision and "image/jpeg" in data and len(data["image/jpeg"])>0:
                        content_image.append({"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64," + data["image/jpeg"] } })
                    elif use_vision and "image/gif" in data and len(data["image/gif"])>0:
                        content_image.append({"type": "image_url", "image_url": {"url": f"data:image/gif;base64," + data["image/gif"] } })
                    elif use_vision and "image/webp" in data and len(data["image/webp"])>0:
                        content_image.append({"type": "image_url", "image_url": {"url": f"data:image/webp;base64," + data["image/webp"] } })
                    # 目前openai vision不支持bmp格式
        
        # 如果是raw单元格，目前暂时没有特殊处理
        if cell["cell_type"] == "raw":
            pass
        
        # 如果是code单元格，目前要处理outputs中的内容
        if cell["cell_type"] == "code":
            if "outputs" in cell and len(cell["outputs"])>0:
                
                for output in cell["outputs"]:
                    # 一般是打印出来的内容
                    if output["output_type"] == "stream":
                        clean_stream = remove_ansi_codes(output["text"])
                        output_text.append(clean_stream.strip())
                    # 单元格输出的错误内容
                    elif output["output_type"] == "error":
                        # 去掉traceback中的颜色类的ansi符号
                        clean_traceback = [remove_ansi_codes(text) for text in output["traceback"]]
                        clean_traceback_text = "\n".join(clean_traceback).strip()
                        output_text.append(f'''Error Name:{output["ename"]}\nError Value:{output["evalue"]}\nError Traceback:{clean_traceback_text}''')
                    elif output["output_type"] == "execute_result" or output["output_type"] == "display_data":
                        if "data" in output and len(output["data"])>0:
                            # 一般是变量输出的值
                            if "text/plain" in output["data"] and len(output["data"]["text/plain"])>0:
                                output_text.append(output["data"]["text/plain"].strip())
                            # 一般是plotly的微缩图片
                            if use_vision and "image/png" in output["data"] and len(output["data"]["image/png"])>0:
                                content_image.append({"type": "image_url", "image_url": {"url": f"data:image/png;base64," + output["data"]["image/png"] } })
                            elif use_vision and "image/jpeg" in output["data"] and len(output["data"]["image/jpeg"])>0:
                                content_image.append({"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64," + output["data"]["image/jpeg"] } })
                            elif use_vision and "image/gif" in output["data"] and len(output["data"]["image/gif"])>0:
                                content_image.append({"type": "image_url", "image_url": {"url": f"data:image/gif;base64," + output["data"]["image/gif"] } })
                            elif use_vision and "image/webp" in output["data"] and len(output["data"]["image/webp"])>0:
                                content_image.append({"type": "image_url", "image_url": {"url": f"data:image/webp;base64," + output["data"]["image/webp"] } })

        # 如果有图片，则标记为有图片
        if len(content_image) > 0:
            has_image = True
        
        # 准备该条信息的结构数据
        content_text = ""
        if len(source_text) > 0:
            content_text += source_text + "\n"
        if len(output_text) > 0:
            content_text += "\nexecuted outputs:\n" + "\n----------\n".join(output_text) + "\n----------"

        content_text = content_text.strip()

        if len(content_image) > 0 and len(content_text) > 0:
            message["content"] = [{"type": "text", "text": content_text}]
            message["content"].extend(content_image)
        elif len(content_image) > 0 and len(content_text) <= 0:
            message["content"] = content_image
        elif len(content_image) <= 0 and len(content_text) > 0:
            message["content"] = content_text
        else:
            continue

        messages.append(message)
        # 这里只计算文本的token数量，不计算图片的token数量
        tokens.append(get_num_tokens(content_text, model))
        
        # 如果是当前活动单元格，则特别标记，因为有的时候，用户可能会放入下文，如果含有下文，则用户当前活动单元格再重复一次
        if cell["num_id"] == active_cell_index and id < len(cell_json_arr)-1:
            is_active_cell_last = False
            last_message = message.copy()
            active_cell_tokens = tokens[-1]
    
    # 最后检查下活动单元格是不是最后一个，如果不是，则再重复一次
    if not is_active_cell_last:
        messages.append(last_message)
        tokens.append(active_cell_tokens)

    return messages, tokens, has_image

# 将csv、excel文件转化为markdown表格
def get_text_from_data(file_path):
    # 根据文件扩展名读取数据
    if file_path.endswith('.csv'):
        df = pd.read_csv(file_path)
    elif file_path.endswith(('.xlsx', '.xls')):
        df = pd.read_excel(file_path)
    df = df.dropna(how='all', axis=1)  # 移除全为空的列
    df = df.dropna(how='all', axis=0)  # 移除全为空的行
    return df.to_markdown(index=False)

# 将pdf转化为文本
def get_text_from_pdf(file_path):
    doc = fitz.open(file_path)
    text = [page.get_text() for page in doc]
    doc.close()
    return "\n".join(text)

# 将docx转化为文本，并分级标注标题
def get_text_from_docx(file_path):
    doc = Document(file_path)
    text = []
    for para in doc.paragraphs:
        if para.text:  # 确保段落有文本
            # 检查段落的样式名称并映射到Markdown标题
            if para.style.name.startswith('Heading'):
                level = para.style.name.replace('Heading ', '')
                markdown_header = '#' * int(level)  # 根据级别确定#的数量
                text.append(f"{markdown_header} {para.text}")
            else:
                text.append(para.text)
    return "\n".join(text)

# 将pptx转化为文本
def get_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        text_boxes = [shape for shape in slide.shapes if hasattr(shape, 'text') and shape.text]
        for index, shape in enumerate(text_boxes):
            # 假设第一个文本框是主标题
            if index == 0:  # 主标题
                text.append(f"# {shape.text}")
            else:  # 其他文本按正常文本处理
                text.append(shape.text)
    return "\n".join(text)

# 将html转化为文本
def get_text_from_html(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        soup = BeautifulSoup(file, 'html5lib')
    return soup.get_text()

# 移除字符串中的ansi颜色代码
def remove_ansi_codes(text):
    ansi_escape = re.compile(r'\x1B[@-_][0-?]*[ -/]*[@-~]')
    return ansi_escape.sub("", text)

# 移除文本中开头和结尾特定文本
def process_source(source, ai_name, user_name, ref_name):
    # 构建正则表达式
    # 匹配以 ai_name 或 user_name 开头的文本
    name_regex = re.compile(r'^(?:{}|{})'.format(re.escape(ai_name), re.escape(user_name)), re.IGNORECASE)

    # 匹配特定格式的 div 标签
    div_pattern = r'<div.*?>{}.*?{}.*?</div>$'.format(re.escape(ref_name), re.escape(ref_name))
    div_regex = re.compile(div_pattern, re.IGNORECASE | re.DOTALL)

    # 移除匹配的文本
    source = name_regex.sub("", source)
    source = div_regex.sub("", source)

    return source

# 计算token数量
def get_num_tokens(text, model):
    # 遇到tiktoken暂时还没有支持的model，使用gpt-4-turbo暂时替代
    try:
        encoding = tiktoken.encoding_for_model(model)
    except KeyError:
        encoding = tiktoken.encoding_for_model("gpt-4-turbo")
    return len(encoding.encode(text))

# 解析参数
def parse_param(data, key, type, default):
    if key in data:
        value = data[key]
        # 如果值已经是目标类型，直接返回
        if isinstance(value, type):
            return value
        # 如果值是字符串，转换为目标类型
        elif isinstance(value, str):
            try:
                # bool类型单独处理，因为可能遇到大小写问题
                if type == bool:
                    value_lower = value.lower()
                    if value_lower == 'true':
                        return True
                    elif value_lower == 'false':
                        return False
                else:
                    return type(value)
            except:
                return default
    # 如果值是其他类型，返回默认值
    return default

def setup_handlers(web_app):
    host_pattern = ".*$"
    base_url = web_app.settings["base_url"]

    # 定义example端点的路由
    get_example_route_pattern = url_path_join(base_url, "jupyterlab-notechat", "get-example")
    get_example_handler = (get_example_route_pattern, RouteHandler)

    # 定义chat端点的路由
    chat_route_pattern = url_path_join(base_url, "jupyterlab-notechat", "chat")
    chat_handler = (chat_route_pattern, ChatHandler)

    # 添加handlers到web应用
    handlers = [get_example_handler, chat_handler]
    web_app.add_handlers(host_pattern, handlers)
